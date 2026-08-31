from __future__ import annotations

from pathlib import Path

import pytest

from agents.doc_parser_agent import DocParserAgent
from orchestrator.graph import build_workflows
from services.assets import AssetStore
from services.catalog import CatalogService
from services.ingestion import IngestionService
from tests.test_ingestion_cdc import EmptyExtractor, MemoryGraph, MemoryVector
from tests.test_orchestrator import FakeQA, FakeUpdater


def _make_documents(root: Path) -> list[Path]:
    from docx import Document
    from openpyxl import Workbook
    from PIL import Image
    from pptx import Presentation
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen.canvas import Canvas

    paths: list[Path] = []
    pdf = root / "sample.pdf"
    canvas = Canvas(str(pdf), pagesize=A4)
    canvas.drawString(40, A4[1] - 40, "Atlas PDF contains searchable project documentation.")
    canvas.save()
    paths.append(pdf)

    docx = root / "sample.docx"
    word = Document()
    word.add_heading("Atlas", level=1)
    word.add_paragraph("Alice owns Atlas.")
    word.save(docx)
    paths.append(docx)

    xlsx = root / "sample.xlsx"
    workbook = Workbook()
    workbook.active.append(["Metric", "Value"])
    workbook.active.append(["Revenue", 100])
    workbook.save(xlsx)
    paths.append(xlsx)

    csv = root / "sample.csv"
    csv.write_text("Metric,Value\nRevenue,100\n", encoding="utf-8")
    paths.append(csv)

    image = root / "sample.png"
    Image.new("RGB", (160, 80), "white").save(image)
    paths.append(image)

    text = root / "sample.txt"
    text.write_text("Atlas plain text", encoding="utf-8")
    paths.append(text)

    markdown = root / "sample.md"
    markdown.write_text("# Atlas\nMarkdown knowledge", encoding="utf-8")
    paths.append(markdown)

    pptx = root / "sample.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "Atlas presentation"
    deck.save(pptx)
    paths.append(pptx)

    html = root / "sample.html"
    html.write_text("<h1>Atlas</h1><p>HTML knowledge</p>", encoding="utf-8")
    paths.append(html)

    json = root / "sample.json"
    json.write_text('{"project":"Atlas","owner":"Alice"}', encoding="utf-8")
    paths.append(json)

    xml = root / "sample.xml"
    xml.write_text("<root><project>Atlas</project></root>", encoding="utf-8")
    paths.append(xml)
    return paths


@pytest.mark.asyncio
async def test_all_eleven_format_families_complete_multinode_ingest(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    paths = _make_documents(tmp_path)
    assert len(paths) == 11
    monkeypatch.setattr(DocParserAgent, "_ocr_image", staticmethod(lambda _image: "image OCR"))
    catalog = CatalogService(str(tmp_path / "catalog.sqlite3"))
    await catalog.init()
    parser = DocParserAgent(asset_store=AssetStore(str(tmp_path / "assets")))
    ingestion = IngestionService(
        parser,
        EmptyExtractor(),
        MemoryVector(),
        MemoryGraph(),
        catalog,
    )
    workflow = build_workflows(ingestion, FakeQA(), FakeUpdater())["ingest"]  # type: ignore[arg-type]
    results = []
    for index, path in enumerate(paths):
        state = await workflow.ainvoke(
            {
                "file_paths": [str(path)],
                "doc_ids": [f"format-{index}"],
                "mime_types": ["application/octet-stream"],
            }
        )
        assert not state.get("error")
        assert "ingest:verify:complete" in state["trace"]
        results.append(state["results"][0])
    assert all(result.status == "READY" and result.chunks_total >= 1 for result in results)
    assert sum(result.modality_counts.get("image", 0) for result in results) >= 1
    assert sum(result.modality_counts.get("table", 0) for result in results) >= 2
    await catalog.close()
