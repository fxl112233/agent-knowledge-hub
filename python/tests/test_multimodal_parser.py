from __future__ import annotations

from pathlib import Path

import pytest

from agents.doc_parser_agent import DocParserAgent, DocumentParseError
from services.assets import AssetStore


@pytest.mark.asyncio
async def test_pptx_preserves_slide_table_and_managed_image(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    image_path = tmp_path / "chart.png"
    Image.new("RGB", (120, 80), "white").save(image_path)
    path = tmp_path / "deck.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "Quarterly report"
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(1.5), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Revenue"
    table.cell(1, 1).text = "100"
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(3), width=Inches(1))
    deck.save(path)

    monkeypatch.setattr(DocParserAgent, "_ocr_image", staticmethod(lambda _image: "chart OCR"))
    asset_store = AssetStore(str(tmp_path / "assets"))
    chunks = await DocParserAgent(asset_store=asset_store).parse(str(path), doc_id="pptx-doc")
    assert any(chunk.metadata.get("slide") == 1 for chunk in chunks)
    assert any(chunk.metadata.get("modality") == "table" for chunk in chunks)
    image = next(chunk for chunk in chunks if chunk.metadata.get("modality") == "image")
    assert Path(image.metadata["asset_path"]).is_file()
    assert asset_store.is_managed(image.metadata["asset_path"])


@pytest.mark.asyncio
async def test_html_json_jsonl_and_xml_keep_structured_locations(tmp_path: Path) -> None:
    html = tmp_path / "page.html"
    html.write_text(
        "<h1>Atlas</h1><p>Owner Alice</p>"
        "<script>remote evil</script><img src='https://example.invalid/tracker.png'>"
        "<table><tr><th>Metric</th><th>Value</th></tr><tr><td>Revenue</td><td>100</td></tr></table>",
        encoding="utf-8",
    )
    json_path = tmp_path / "data.json"
    json_path.write_text('{"project":{"owner":"Alice","amount":100}}', encoding="utf-8")
    jsonl = tmp_path / "events.jsonl"
    jsonl.write_text('{"event":"created"}\n{"event":"updated"}\n', encoding="utf-8")
    xml = tmp_path / "data.xml"
    xml.write_text('<root><project id="1"><owner>Alice</owner></project></root>', encoding="utf-8")

    parser = DocParserAgent(asset_store=AssetStore(str(tmp_path / "assets")))
    html_chunks = await parser.parse(str(html), doc_id="html")
    json_chunks = await parser.parse(str(json_path), doc_id="json")
    jsonl_chunks = await parser.parse(str(jsonl), doc_id="jsonl")
    xml_chunks = await parser.parse(str(xml), doc_id="xml")

    assert any(chunk.metadata.get("section") == "h1:Atlas" for chunk in html_chunks)
    assert any(chunk.metadata.get("modality") == "table" for chunk in html_chunks)
    assert all("remote evil" not in chunk.content for chunk in html_chunks)
    assert any("['owner']" in chunk.metadata.get("json_path", "") for chunk in json_chunks)
    assert {chunk.metadata["json_path"].split("]", 1)[0] + "]" for chunk in jsonl_chunks} == {
        "$[1]",
        "$[2]",
    }
    assert any("/owner[1]" in chunk.metadata.get("xpath", "") for chunk in xml_chunks)


@pytest.mark.asyncio
async def test_structured_limits_xxe_and_corrupt_new_formats(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    parser = DocParserAgent(asset_store=AssetStore(str(tmp_path / "assets")))
    deep = tmp_path / "deep.json"
    deep.write_text('{"a":{"b":{"c":{"d":1}}}}', encoding="utf-8")
    monkeypatch.setattr("agents.doc_parser_agent.settings.max_json_depth", 3)
    with pytest.raises(DocumentParseError, match="depth"):
        await parser.parse(str(deep), doc_id="deep")

    xxe = tmp_path / "xxe.xml"
    xxe.write_text(
        '<!DOCTYPE root [<!ENTITY xxe SYSTEM "file:///etc/passwd">]><root>&xxe;</root>',
        encoding="utf-8",
    )
    with pytest.raises(DocumentParseError, match="unsafe or invalid XML"):
        await parser.parse(str(xxe), doc_id="xxe")

    for name in ("broken.pptx", "broken.json", "broken.jsonl", "broken.xml"):
        path = tmp_path / name
        path.write_text("not valid", encoding="utf-8")
        with pytest.raises(DocumentParseError):
            await parser.parse(str(path), doc_id=path.stem)


@pytest.mark.asyncio
async def test_vision_failure_falls_back_to_ocr_and_keeps_asset(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    from PIL import Image

    class FailingVision:
        async def describe(self, _image):
            raise TimeoutError("vision timeout")

    path = tmp_path / "scan.png"
    Image.new("RGB", (100, 60), "white").save(path)
    monkeypatch.setattr(DocParserAgent, "_ocr_image", staticmethod(lambda _image: "OCR survives"))
    parser = DocParserAgent(vision_model=FailingVision(), asset_store=AssetStore(str(tmp_path / "assets")))
    chunks = await parser.parse(str(path), doc_id="scan")
    assert chunks[0].content == "OCR survives"
    assert chunks[0].metadata["vision_fallback"] == "TimeoutError"
    assert Path(chunks[0].metadata["asset_path"]).is_file()


def test_asset_store_deduplicates_prunes_and_deletes_document(tmp_path: Path) -> None:
    from PIL import Image

    store = AssetStore(str(tmp_path / "assets"))
    first = store.save_image("doc", Image.new("RGB", (10, 10), "red"))
    duplicate = store.save_image("doc", Image.new("RGB", (10, 10), "red"))
    second = store.save_image("doc", Image.new("RGB", (10, 10), "blue"))
    assert first == duplicate
    assert store.prune_document("doc", [first]) == 1
    assert Path(first).exists() and not Path(second).exists()
    assert store.delete_document("doc") == 1
    assert not Path(first).exists()


def test_asset_store_maps_external_document_ids_to_safe_directories(tmp_path: Path) -> None:
    from PIL import Image

    store = AssetStore(str(tmp_path / "assets"))
    external_id = "https://example.com/reports/annual?id=42"
    asset = Path(store.save_image(external_id, Image.new("RGB", (10, 10), "red")))

    assert asset.is_file()
    assert asset.parent.name.startswith("external-")
    assert asset.parent.parent == (tmp_path / "assets").resolve()
    assert store.delete_document(external_id) == 1
    assert not asset.exists()
