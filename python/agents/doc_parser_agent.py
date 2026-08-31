"""Safe multimodal parser for the eleven supported document format families."""

from __future__ import annotations

import asyncio
import hashlib
import io
import json
import os
import re
from collections import defaultdict
from dataclasses import dataclass, field
from enum import StrEnum
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

from config import settings
from services.assets import AssetStore


class DocumentParseError(ValueError):
    """A supported document could not be parsed."""


class UnsupportedDocumentError(DocumentParseError):
    """The file extension is not supported."""


class _SafeHTMLCollector(HTMLParser):
    """Extract visible text and tables; no network-capable HTML features are used."""

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.units: list[ParsedUnit] = []
        self._heading = "document"
        self._text: list[str] = []
        self._table_rows: list[list[str]] = []
        self._row: list[str] = []
        self._cell: list[str] = []
        self._in_table = False
        self._in_cell = False
        self._heading_tag = ""
        self._heading_buffer: list[str] = []
        self._skip_depth = 0
        self._section_index = 0
        self._table_index = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        del attrs
        tag = tag.lower()
        if tag in {"script", "style", "noscript", "template", "svg"}:
            self._skip_depth += 1
        elif self._skip_depth:
            return
        elif tag in {"h1", "h2", "h3", "h4", "h5", "h6"}:
            self._flush_text()
            self._heading_tag = tag
            self._heading_buffer = []
        elif tag == "table":
            self._flush_text()
            self._in_table = True
            self._table_rows = []
        elif tag == "tr" and self._in_table:
            self._row = []
        elif tag in {"td", "th"} and self._in_table:
            self._in_cell = True
            self._cell = []

    def handle_endtag(self, tag: str) -> None:
        tag = tag.lower()
        if tag in {"script", "style", "noscript", "template", "svg"} and self._skip_depth:
            self._skip_depth -= 1
            return
        if self._skip_depth:
            return
        if tag == self._heading_tag:
            heading = " ".join(self._heading_buffer).strip()
            if heading:
                self._heading = f"{tag}:{heading}"
                self._text.append(heading)
            self._heading_tag = ""
            self._heading_buffer = []
        elif tag in {"td", "th"} and self._in_cell:
            self._row.append(" ".join(self._cell).strip())
            self._in_cell = False
        elif tag == "tr" and self._in_table and self._row:
            self._table_rows.append(self._row)
        elif tag == "table" and self._in_table:
            self._table_index += 1
            content = "\n".join(" | ".join(row) for row in self._table_rows).strip()
            if content:
                self.units.append(
                    ParsedUnit(
                        content,
                        f"table:{self._table_index}",
                        {"table": self._table_index, "kind": "table", "modality": "table"},
                    )
                )
            self._in_table = False

    def handle_data(self, data: str) -> None:
        if self._skip_depth or not data.strip():
            return
        if self._heading_tag:
            self._heading_buffer.append(data.strip())
        elif self._in_cell:
            self._cell.append(data.strip())
        elif not self._in_table:
            self._text.append(data.strip())

    def close(self) -> None:
        super().close()
        self._flush_text()

    def _flush_text(self) -> None:
        content = "\n".join(self._text).strip()
        if content:
            self._section_index += 1
            self.units.append(
                ParsedUnit(
                    content,
                    f"section:{self._section_index}",
                    {"section": self._heading, "kind": "text", "modality": "text"},
                )
            )
        self._text.clear()


class DocType(StrEnum):
    PDF = "pdf"
    WORD = "word"
    IMAGE = "image"
    TABLE = "table"
    TEXT = "text"
    MARKDOWN = "markdown"
    PRESENTATION = "presentation"
    HTML = "html"
    JSON = "json"
    XML = "xml"


@dataclass(slots=True)
class ParsedUnit:
    content: str
    unit_id: str
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass(slots=True)
class DocumentChunk:
    content: str
    doc_id: str
    chunk_index: int
    doc_type: DocType
    metadata: dict[str, Any] = field(default_factory=dict)
    embedding: list[float] | None = None
    identifier: str = ""
    content_hash: str = ""

    @property
    def chunk_id(self) -> str:
        return self.identifier or f"{self.doc_id}#chunk-{self.chunk_index}"


class DocParserAgent:
    """Parse supported files into token-aware chunks with source metadata."""

    SUPPORTED_EXTENSIONS: dict[str, DocType] = {
        ".pdf": DocType.PDF,
        ".docx": DocType.WORD,
        ".png": DocType.IMAGE,
        ".jpg": DocType.IMAGE,
        ".jpeg": DocType.IMAGE,
        ".csv": DocType.TABLE,
        ".xlsx": DocType.TABLE,
        ".txt": DocType.TEXT,
        ".md": DocType.MARKDOWN,
        ".pptx": DocType.PRESENTATION,
        ".html": DocType.HTML,
        ".htm": DocType.HTML,
        ".json": DocType.JSON,
        ".jsonl": DocType.JSON,
        ".xml": DocType.XML,
    }

    def __init__(self, vision_model: Any | None = None, asset_store: AssetStore | None = None) -> None:
        self.vision_model = vision_model
        self.asset_store = asset_store or AssetStore()

    async def parse(
        self,
        file_path: str,
        *,
        doc_id: str | None = None,
        version: int = 1,
    ) -> list[DocumentChunk]:
        path = Path(file_path).resolve()
        if not path.is_file():
            raise DocumentParseError(f"file does not exist: {path}")
        doc_type = self._classify(path)
        doc_id = doc_id or self._make_doc_id(str(path))

        if doc_type == DocType.PDF:
            units = await self._parse_pdf(path, doc_id)
        elif doc_type == DocType.WORD:
            units = await self._parse_docx(path, doc_id)
        elif doc_type == DocType.IMAGE:
            units = await self._parse_image(path, doc_id)
        elif doc_type == DocType.TABLE:
            units = self._parse_csv(path) if path.suffix.lower() == ".csv" else self._parse_excel(path)
        elif doc_type == DocType.MARKDOWN:
            units = self._parse_markdown(path)
        elif doc_type == DocType.PRESENTATION:
            units = await self._parse_pptx(path, doc_id)
        elif doc_type == DocType.HTML:
            units = self._parse_html(path)
        elif doc_type == DocType.JSON:
            units = self._parse_json(path)
        elif doc_type == DocType.XML:
            units = self._parse_xml(path)
        else:
            units = self._parse_text(path)

        units = [unit for unit in units if unit.content.strip()]
        if not units:
            raise DocumentParseError(f"no readable content found in {path.name}")
        return self._chunk_units(units, doc_id, doc_type, str(path), version)

    async def parse_batch(
        self,
        file_paths: list[str],
        *,
        doc_ids: list[str] | None = None,
        version: int = 1,
    ) -> list[DocumentChunk]:
        if doc_ids is not None and len(doc_ids) != len(file_paths):
            raise ValueError("doc_ids must match file_paths")
        tasks = [
            self.parse(path, doc_id=doc_ids[index] if doc_ids else None, version=version)
            for index, path in enumerate(file_paths)
        ]
        parsed = await asyncio.gather(*tasks)
        return [chunk for chunks in parsed for chunk in chunks]

    def _classify(self, path: Path) -> DocType:
        try:
            return self.SUPPORTED_EXTENSIONS[path.suffix.lower()]
        except KeyError as exc:
            raise UnsupportedDocumentError(f"unsupported extension: {path.suffix or '<none>'}") from exc

    @staticmethod
    def _make_doc_id(file_path: str) -> str:
        return hashlib.sha256(str(Path(file_path).resolve()).encode("utf-8")).hexdigest()[:20]

    async def _parse_pdf(self, path: Path, doc_id: str) -> list[ParsedUnit]:
        try:
            import pymupdf
        except ImportError as exc:
            raise DocumentParseError("PyMuPDF is required for PDF parsing") from exc

        units: list[ParsedUnit] = []
        try:
            document = pymupdf.open(path)
        except Exception as exc:
            raise DocumentParseError(f"invalid PDF: {path.name}") from exc
        try:
            if document.page_count > settings.max_pdf_pages:
                raise DocumentParseError(
                    f"PDF contains {document.page_count} pages; limit is {settings.max_pdf_pages}"
                )
            tables_by_page = self._extract_pdf_tables(path)
            for page_index, page in enumerate(document):
                text = page.get_text("text").strip()
                used_ocr = False
                asset_path = ""
                vision_fallback = ""
                if len(text) < settings.ocr_min_text_chars:
                    pix = page.get_pixmap(matrix=pymupdf.Matrix(2, 2), alpha=False)
                    from PIL import Image

                    image = Image.open(io.BytesIO(pix.tobytes("png")))
                    asset_path = self.asset_store.save_image(doc_id, image)
                    ocr = self._ocr_image(image)
                    if len(ocr.strip()) > len(text):
                        text = ocr.strip()
                        used_ocr = True
                    description, vision_fallback = await self._vision_description(image)
                    text = "\n".join(value for value in (text, description) if value)
                if text:
                    units.append(
                        ParsedUnit(
                            text,
                            f"page:{page_index + 1}:text",
                            {
                                "page": page_index + 1,
                                "kind": "image" if asset_path else "text",
                                "modality": "image" if asset_path else "text",
                                "ocr": used_ocr,
                                "asset_path": asset_path,
                                "vision_fallback": vision_fallback,
                            },
                        )
                    )
                for table_index, table in enumerate(tables_by_page.get(page_index + 1, []), start=1):
                    units.append(
                        ParsedUnit(
                            table,
                            f"page:{page_index + 1}:table:{table_index}",
                            {
                                "page": page_index + 1,
                                "table": table_index,
                                "kind": "table",
                                "modality": "table",
                            },
                        )
                    )
        finally:
            document.close()
        return units

    @staticmethod
    def _extract_pdf_tables(path: Path) -> dict[int, list[str]]:
        tables: dict[int, list[str]] = defaultdict(list)
        try:
            import pdfplumber

            with pdfplumber.open(path) as pdf:
                for page_number, page in enumerate(pdf.pages, start=1):
                    for table in page.extract_tables() or []:
                        rows = [
                            " | ".join("" if cell is None else str(cell) for cell in row) for row in table
                        ]
                        value = "\n".join(rows).strip()
                        if value:
                            tables[page_number].append(value)
        except Exception:
            # Text/OCR remains usable when table detection cannot handle a page.
            return {}
        return dict(tables)

    async def _parse_docx(self, path: Path, doc_id: str) -> list[ParsedUnit]:
        try:
            from docx import Document
        except ImportError as exc:
            raise DocumentParseError("python-docx is required for DOCX parsing") from exc
        try:
            document = Document(path)
        except Exception as exc:
            raise DocumentParseError(f"invalid DOCX: {path.name}") from exc

        units: list[ParsedUnit] = []
        section = "document"
        buffer: list[str] = []

        def flush() -> None:
            if buffer:
                units.append(
                    ParsedUnit(
                        "\n".join(buffer),
                        f"section:{len(units) + 1}",
                        {"section": section, "kind": "text", "modality": "text"},
                    )
                )
                buffer.clear()

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            if paragraph.style and paragraph.style.name.lower().startswith("heading"):
                flush()
                section = text
            buffer.append(text)
        flush()

        for table_index, table in enumerate(document.tables, start=1):
            rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
            value = "\n".join(rows).strip()
            if value:
                units.append(
                    ParsedUnit(
                        value,
                        f"table:{table_index}",
                        {"table": table_index, "kind": "table", "modality": "table"},
                    )
                )

        image_index = 0
        for relationship in document.part.rels.values():
            if "image" not in relationship.reltype:
                continue
            image_index += 1
            try:
                from PIL import Image

                image = Image.open(io.BytesIO(relationship.target_part.blob))
                text = self._ocr_image(image)
                description, fallback = await self._vision_description(image)
                asset_path = self.asset_store.save_image(doc_id, image)
                content = "\n".join(value.strip() for value in (text, description) if value.strip())
                units.append(
                    ParsedUnit(
                        content or "[Image without extracted text]",
                        f"image:{image_index}",
                        {
                            "image": image_index,
                            "kind": "image",
                            "modality": "image",
                            "asset_path": asset_path,
                            "ocr": bool(text.strip()),
                            "vision_fallback": fallback,
                        },
                    )
                )
            except Exception:
                continue
        return units

    async def _parse_image(self, path: Path, doc_id: str) -> list[ParsedUnit]:
        try:
            from PIL import Image

            image = Image.open(path)
            image.load()
        except Exception as exc:
            raise DocumentParseError(f"invalid image: {path.name}") from exc
        text = self._ocr_image(image)
        description, fallback = await self._vision_description(image)
        asset_path = self.asset_store.save_image(doc_id, image)
        content = "\n".join(value for value in (text.strip(), description.strip()) if value)
        return [
            ParsedUnit(
                content or "[Image without extracted text]",
                "image:1",
                {
                    "image": 1,
                    "kind": "image",
                    "modality": "image",
                    "ocr": bool(text.strip()),
                    "asset_path": asset_path,
                    "vision_fallback": fallback,
                },
            )
        ]

    @staticmethod
    def _ocr_image(image: Any) -> str:
        try:
            import pytesseract

            # Uniform document blocks are more stable with PSM 6, especially after
            # small rotations where automatic layout detection may reorder lines.
            return pytesseract.image_to_string(
                image,
                lang=settings.ocr_languages,
                config="--psm 6",
            )
        except Exception:
            return ""

    async def _vision_description(self, image: Any) -> tuple[str, str]:
        if self.vision_model is None or not settings.vision_enabled:
            return "", "disabled_or_unconfigured"
        try:
            return await self._describe_image(image), ""
        except Exception as exc:
            return "", type(exc).__name__

    async def _describe_image(self, image: Any) -> str:
        describe = getattr(self.vision_model, "describe", None)
        if describe is not None:
            return str(await describe(image))
        import base64

        buffer = io.BytesIO()
        image.convert("RGB").save(buffer, format="JPEG", quality=85)
        encoded = base64.b64encode(buffer.getvalue()).decode("ascii")
        response = await self.vision_model.ainvoke(
            [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "请客观描述文档图片中的文字、表格、图表、数值及其关系。",
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{encoded}",
                                "detail": "low",
                            },
                        },
                    ],
                }
            ]
        )
        return str(getattr(response, "content", response))

    @staticmethod
    def _parse_csv(path: Path) -> list[ParsedUnit]:
        import csv

        from charset_normalizer import from_bytes

        raw = path.read_bytes()
        match = from_bytes(raw).best()
        text = str(match) if match is not None else raw.decode("utf-8-sig")
        rows = list(csv.reader(io.StringIO(text)))
        if len(rows) > settings.max_spreadsheet_rows:
            raise DocumentParseError(f"CSV row limit exceeded: {len(rows)}")
        units: list[ParsedUnit] = []
        for start in range(0, len(rows), 20):
            batch = rows[start : start + 20]
            value = "\n".join(" | ".join(str(cell) for cell in row) for row in batch).strip()
            if value:
                units.append(
                    ParsedUnit(
                        value,
                        f"rows:{start + 1}-{start + len(batch)}",
                        {
                            "row_start": start + 1,
                            "row_end": start + len(batch),
                            "kind": "table",
                            "modality": "table",
                        },
                    )
                )
        return units

    @staticmethod
    def _parse_excel(path: Path) -> list[ParsedUnit]:
        try:
            import openpyxl

            workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
        except Exception as exc:
            raise DocumentParseError(f"invalid XLSX: {path.name}") from exc
        units: list[ParsedUnit] = []
        total_rows = 0
        try:
            for sheet in workbook.worksheets:
                rows: list[list[str]] = []
                for row in sheet.iter_rows(values_only=True):
                    total_rows += 1
                    if total_rows > settings.max_spreadsheet_rows:
                        raise DocumentParseError("spreadsheet row limit exceeded")
                    rows.append(["" if cell is None else str(cell) for cell in row])
                    if len(rows) == 20:
                        units.append(DocParserAgent._sheet_unit(sheet.title, total_rows - 19, rows))
                        rows = []
                if rows:
                    units.append(DocParserAgent._sheet_unit(sheet.title, total_rows - len(rows) + 1, rows))
        finally:
            workbook.close()
        return units

    @staticmethod
    def _sheet_unit(sheet: str, row_start: int, rows: list[list[str]]) -> ParsedUnit:
        row_end = row_start + len(rows) - 1
        value = "\n".join(" | ".join(row) for row in rows)
        return ParsedUnit(
            value,
            f"sheet:{sheet}:rows:{row_start}-{row_end}",
            {
                "sheet": sheet,
                "row_start": row_start,
                "row_end": row_end,
                "kind": "table",
                "modality": "table",
            },
        )

    @staticmethod
    def _parse_text(path: Path) -> list[ParsedUnit]:
        text = DocParserAgent._read_text(path)
        return [ParsedUnit(text, "document:1", {"kind": "text", "modality": "text"})]

    @staticmethod
    def _read_text(path: Path) -> str:
        from charset_normalizer import from_bytes

        raw = path.read_bytes()
        match = from_bytes(raw).best()
        return str(match) if match is not None else raw.decode("utf-8-sig")

    @staticmethod
    def _parse_markdown(path: Path) -> list[ParsedUnit]:
        """Split Markdown at headings so unrelated sections keep stable IDs."""
        text = DocParserAgent._read_text(path)
        units: list[ParsedUnit] = []
        heading = "preamble"
        buffer: list[str] = []
        occurrences: defaultdict[str, int] = defaultdict(int)

        def flush() -> None:
            content = "\n".join(buffer).strip()
            if not content:
                buffer.clear()
                return
            normalized = DocParserAgent._normalize(heading).casefold()
            occurrence = occurrences[normalized]
            occurrences[normalized] += 1
            section_key = hashlib.sha256(normalized.encode("utf-8")).hexdigest()[:16]
            units.append(
                ParsedUnit(
                    content,
                    f"section:{section_key}:{occurrence}",
                    {"section": heading, "kind": "text", "modality": "text"},
                )
            )
            buffer.clear()

        for line in text.splitlines():
            match = re.match(r"^\s{0,3}#{1,6}\s+(.+?)\s*#*\s*$", line)
            if match:
                flush()
                heading = match.group(1).strip()
            buffer.append(line)
        flush()
        return units

    async def _parse_pptx(self, path: Path, doc_id: str) -> list[ParsedUnit]:
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            presentation = Presentation(path)
        except Exception as exc:
            raise DocumentParseError(f"invalid PPTX: {path.name}") from exc
        units: list[ParsedUnit] = []
        image_index = 0
        for slide_number, slide in enumerate(presentation.slides, start=1):
            title = ""
            if slide.shapes.title is not None:
                title = slide.shapes.title.text.strip()
            for shape_index, shape in enumerate(slide.shapes, start=1):
                unit_prefix = f"slide:{slide_number}:shape:{shape_index}"
                if getattr(shape, "has_table", False):
                    rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in shape.table.rows]
                    value = "\n".join(rows).strip()
                    if value:
                        units.append(
                            ParsedUnit(
                                value,
                                f"{unit_prefix}:table",
                                {
                                    "slide": slide_number,
                                    "title": title,
                                    "table": shape_index,
                                    "kind": "table",
                                    "modality": "table",
                                },
                            )
                        )
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_index += 1
                    try:
                        from PIL import Image

                        image = Image.open(io.BytesIO(shape.image.blob))
                        image.load()
                        ocr = self._ocr_image(image)
                        description, fallback = await self._vision_description(image)
                        asset_path = self.asset_store.save_image(doc_id, image)
                        content = "\n".join(value.strip() for value in (ocr, description) if value.strip())
                        units.append(
                            ParsedUnit(
                                content or "[Image without extracted text]",
                                f"{unit_prefix}:image",
                                {
                                    "slide": slide_number,
                                    "title": title,
                                    "image": image_index,
                                    "kind": "image",
                                    "modality": "image",
                                    "asset_path": asset_path,
                                    "ocr": bool(ocr.strip()),
                                    "vision_fallback": fallback,
                                },
                            )
                        )
                    except Exception:
                        continue
                elif getattr(shape, "has_text_frame", False):
                    value = "\n".join(
                        paragraph.text.strip()
                        for paragraph in shape.text_frame.paragraphs
                        if paragraph.text.strip()
                    )
                    if value:
                        units.append(
                            ParsedUnit(
                                value,
                                f"{unit_prefix}:text",
                                {
                                    "slide": slide_number,
                                    "title": title,
                                    "text_box": shape_index,
                                    "kind": "text",
                                    "modality": "text",
                                },
                            )
                        )
        return units

    @staticmethod
    def _parse_html(path: Path) -> list[ParsedUnit]:
        parser = _SafeHTMLCollector()
        try:
            parser.feed(DocParserAgent._read_text(path))
            parser.close()
        except Exception as exc:
            raise DocumentParseError(f"invalid HTML: {path.name}") from exc
        return parser.units

    @staticmethod
    def _parse_json(path: Path) -> list[ParsedUnit]:
        try:
            if path.suffix.lower() == ".jsonl":
                roots = [
                    (f"$[{line_number}]", json.loads(line))
                    for line_number, line in enumerate(DocParserAgent._read_text(path).splitlines(), start=1)
                    if line.strip()
                ]
            else:
                roots = [("$", json.loads(DocParserAgent._read_text(path)))]
        except (UnicodeError, json.JSONDecodeError) as exc:
            raise DocumentParseError(f"invalid JSON: {path.name}") from exc

        units: list[ParsedUnit] = []
        node_count = 0
        stack: list[tuple[str, Any, int]] = [(json_path, value, 1) for json_path, value in reversed(roots)]
        while stack:
            json_path, value, depth = stack.pop()
            node_count += 1
            if node_count > settings.max_structured_nodes:
                raise DocumentParseError("JSON node limit exceeded")
            if depth > settings.max_json_depth:
                raise DocumentParseError("JSON depth limit exceeded")
            if isinstance(value, dict):
                for key, child in reversed(list(value.items())):
                    escaped = str(key).replace("\\", "\\\\").replace("'", "\\'")
                    stack.append((f"{json_path}['{escaped}']", child, depth + 1))
            elif isinstance(value, list):
                for index in range(len(value) - 1, -1, -1):
                    stack.append((f"{json_path}[{index}]", value[index], depth + 1))
            else:
                rendered = json.dumps(value, ensure_ascii=False)
                units.append(
                    ParsedUnit(
                        f"{json_path} = {rendered}",
                        f"json:{hashlib.sha256(json_path.encode()).hexdigest()[:20]}",
                        {"json_path": json_path, "kind": "text", "modality": "text"},
                    )
                )
        return units

    @staticmethod
    def _parse_xml(path: Path) -> list[ParsedUnit]:
        try:
            from defusedxml import ElementTree

            root = ElementTree.parse(path).getroot()
        except Exception as exc:
            raise DocumentParseError(f"unsafe or invalid XML: {path.name}") from exc

        units: list[ParsedUnit] = []
        node_count = 0
        stack: list[tuple[Any, str, int]] = [(root, f"/{root.tag}[1]", 1)]
        while stack:
            element, xpath, depth = stack.pop()
            node_count += 1
            if node_count > settings.max_structured_nodes:
                raise DocumentParseError("XML node limit exceeded")
            if depth > settings.max_json_depth:
                raise DocumentParseError("XML depth limit exceeded")
            values = [f"@{name}={value}" for name, value in sorted(element.attrib.items())]
            if element.text and element.text.strip():
                values.append(element.text.strip())
            if values:
                units.append(
                    ParsedUnit(
                        f"{xpath} " + " | ".join(values),
                        f"xml:{hashlib.sha256(xpath.encode()).hexdigest()[:20]}",
                        {"xpath": xpath, "kind": "text", "modality": "text"},
                    )
                )
            children = list(element)
            totals: defaultdict[str, int] = defaultdict(int)
            indexed: list[tuple[Any, str]] = []
            for child in children:
                totals[str(child.tag)] += 1
                indexed.append((child, f"{xpath}/{child.tag}[{totals[str(child.tag)]}]"))
            for child, child_xpath in reversed(indexed):
                stack.append((child, child_xpath, depth + 1))
        return units

    def _chunk_units(
        self,
        units: list[ParsedUnit],
        doc_id: str,
        doc_type: DocType,
        source: str,
        version: int,
    ) -> list[DocumentChunk]:
        chunks: list[DocumentChunk] = []
        occurrences: defaultdict[tuple[str, str], int] = defaultdict(int)
        for unit in units:
            for content, token_start, token_end in self._token_windows(unit.content):
                normalized = self._normalize(content)
                if not normalized:
                    continue
                content_hash = hashlib.sha256(normalized.encode("utf-8")).hexdigest()
                occurrence_key = (unit.unit_id, content_hash)
                occurrence = occurrences[occurrence_key]
                occurrences[occurrence_key] += 1
                identifier = hashlib.sha256(
                    f"{doc_id}|{unit.unit_id}|{content_hash}|{occurrence}".encode()
                ).hexdigest()[:32]
                metadata = {
                    **unit.metadata,
                    "source": source,
                    "file_name": os.path.basename(source),
                    "unit_id": unit.unit_id,
                    "token_start": token_start,
                    "token_end": token_end,
                    "version": version,
                }
                chunks.append(
                    DocumentChunk(
                        content=content.strip(),
                        doc_id=doc_id,
                        chunk_index=len(chunks),
                        doc_type=doc_type,
                        metadata=metadata,
                        identifier=identifier,
                        content_hash=content_hash,
                    )
                )
        return chunks

    def _token_windows(self, text: str) -> list[tuple[str, int, int]]:
        try:
            import tiktoken

            encoding = tiktoken.get_encoding("cl100k_base")
            tokens = encoding.encode(text)
            size = settings.chunk_size_tokens
            overlap = min(settings.chunk_overlap_tokens, size - 1)
            windows: list[tuple[str, int, int]] = []
            start = 0
            while start < len(tokens):
                end = min(start + size, len(tokens))
                windows.append((encoding.decode(tokens[start:end]), start, end))
                if end == len(tokens):
                    break
                start = end - overlap
            return windows
        except Exception:
            words = re.findall(r"\S+", text)
            if not words:
                return []
            size = settings.chunk_size_tokens
            overlap = min(settings.chunk_overlap_tokens, size - 1)
            windows = []
            start = 0
            while start < len(words):
                end = min(start + size, len(words))
                windows.append((" ".join(words[start:end]), start, end))
                if end == len(words):
                    break
                start = end - overlap
            return windows

    @staticmethod
    def _normalize(text: str) -> str:
        return re.sub(r"\s+", " ", text).strip()
